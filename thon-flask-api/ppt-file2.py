from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_layout_title = prs.slide_layouts[0] # Typically the title slide layout
slide = prs.slides.add_slide(slide_layout_title)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Modernizing Our Architecture: Enabling Agility and Growth"
subtitle.text = "Addressing the Limitations of Our Monolithic System in [Cleaning Services Company Name]\n[Your Name/Team]\n[Date]"

# Slide 2: Challenges of Our Monolithic Architecture
slide_layout_bullet = prs.slide_layouts[1] # Typically Title and Body layout
slide = prs.slides.add_slide(slide_layout_bullet)
title = slide.shapes.title
body_shape = slide.shapes.placeholders[0]
tf = body_shape.text_frame

title.text = "Challenges of Our Monolithic Architecture"

# Add bullet points to the body
p = tf.add_paragraph()
p.text = "Limited Flexibility and Adaptability:"
tf.level = 1
p = tf.add_paragraph()
p.text = "Slow and risky changes, difficulty adopting new tech."
tf.level = 0 # Reset level

p = tf.add_paragraph()
p.text = "Scalability Issues:"
tf.level = 1
p = tf.add_paragraph()
p.text = "Entire application must scale, inefficient resource use."
tf.level = 0

p = tf.add_paragraph()
p.text = "Resilience Concerns:"
tf.level = 1
p = tf.add_paragraph()
p.text = "Single point of failure, difficult fault isolation."
tf.level = 0

p = tf.add_paragraph()
p.text = "Development Bottlenecks:"
tf.level = 1
p = tf.add_paragraph()
p.text = "Complex codebase, lengthy deployments."
tf.level = 0


# Slide 3: Benefits of Microservices Architecture (Comparison)
# Using a blank layout or a title and content layout for a table/comparison list
slide_layout_title_content = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout_title_content)
title = slide.shapes.title
body_shape = slide.shapes.placeholders[0]
tf = body_shape.text_frame

title.text = "Benefits: Microservices vs. Monolith"

# Add comparative bullet points
p = tf.add_paragraph()
p.text = "Development Agility:"
tf.level = 1
p = tf.add_paragraph()
p.text = "Monolith: Slow, risky. Microservices: Faster cycles, independent deployments."
tf.level = 0

p = tf.add_paragraph()
p.text = "Scalability:"
tf.level = 1
p = tf.add_paragraph()
p.text = "Monolith: Scale entire app. Microservices: Scale individual services."
tf.level = 0

p = tf.add_paragraph()
p.text = "Resilience:"
tf.level = 1
p = tf.add_paragraph()
p.text = "Monolith: Single point of failure. Microservices: Enhanced fault isolation."
tf.level = 0

p = tf.add_paragraph()
p.text = "Adaptability & Innovation:"
tf.level = 1
p = tf.add_paragraph()
p.text = "Monolith: Difficult technology adoption. Microservices: Technology diversity, faster innovation."
tf.level = 0

p = tf.add_paragraph()
p.text = "Alignment with Business Goals:"
tf.level = 1
p = tf.add_paragraph()
p.text = "Microservices enable teams to focus on specific business capabilities for better alignment."
tf.level = 0


# Slide 4: Recommendations and Next Steps
slide = prs.slides.add_slide(slide_layout_bullet)
title = slide.shapes.title
body_shape = slide.shapes.placeholders[0]
tf = body_shape.text_frame

title.text = "Recommendations and Next Steps"

p = tf.add_paragraph()
p.text = "Adopt a phased migration approach (e.g., Strangler Fig Pattern)."
tf.level = 0

p = tf.add_paragraph()
p.text = "Identify and extract a low-risk service first."
tf.level = 0

p = tf.add_paragraph()
p.text = "Invest in necessary tools (CI/CD, monitoring) and team training."
tf.level = 0

p = tf.add_paragraph()
p.text = "Foster a DevOps culture."
tf.level = 0

p = tf.add_paragraph()
p.text = "Form a working group to plan the first service extraction."
tf.level = 0


# Save the presentation
prs.save("Microservices_Presentation_[Cleaning Services Company Name].pptx")

print("PowerPoint presentation 'Microservices_Presentation_[Cleaning Services Company Name].pptx' created successfully.")