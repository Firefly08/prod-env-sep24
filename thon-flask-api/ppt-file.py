from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Slide 1: Introduction
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Microservices Architecture for Cleaning Services Company"
content = slide.placeholders[1].text_frame
content.text = "Introduction"
p = content.add_paragraph()
p.text = "Objective: Present findings on transitioning from a monolithic to a microservices architecture."
p.level = 1
p = content.add_paragraph()
p.text = "Context: Current monolithic application hinders adaptability and scalability."
p.level = 1
p = content.add_paragraph()
p.text = "Goal: Propose a microservices-based solution to enhance flexibility and support business objectives."
p.level = 1
# Add an icon (simulated as a shape)
icon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(0.5), Inches(1), Inches(1))
icon.fill.solid()
icon.fill.fore_color.rgb = RGBColor(0, 105, 148)
icon.line.color.rgb = RGBColor(0, 0, 0)

# Slide 2: Analysis of Monolithic Architecture
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Analysis of Monolithic Architecture"
content = slide.placeholders[1].text_frame
content.text = "Details"
p = content.add_paragraph()
p.text = "Application: Centralised booking and management system for cleaning services."
p.level = 1
p = content.add_paragraph()
p.text = "Limitations:"
p.level = 1
p = content.add_paragraph()
p.text = "Code Complexity: Single codebase with intertwined modules (booking, invoicing, scheduling) increases maintenance effort."
p.level = 2
p = content.add_paragraph()
p.text = "Deployment Challenges: Changes require full system redeployment, causing downtime."
p.level = 2
p = content.add_paragraph()
p.text = "Scalability Issues: Scaling the entire application is resource-intensive, even for specific high-demand features."
p.level = 2
p = content.add_paragraph()
p.text = "Resilience: A single failure (e.g., database issue) impacts all functionalities."
p.level = 2
p = content.add_paragraph()
p.text = "Impact: Slow response to market changes, delayed feature rollouts, and reduced customer satisfaction."
p.level = 1
# Add an icon (simulated as a shape)
icon = slide.shapes.add_shape(MSO_SHAPE.TRIANGLE, Inches(8), Inches(0.5), Inches(1), Inches(1))
icon.fill.solid()
icon.fill.fore_color.rgb = RGBColor(255, 165, 0)
icon.line.color.rgb = RGBColor(0, 0, 0)

# Slide 3: Microservices Architecture Overview
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Microservices Architecture Overview"
content = slide.placeholders[1].text_frame
content.text = "Key Points"
p = content.add_paragraph()
p.text = "Key Characteristics:"
p.level = 1
p = content.add_paragraph()
p.text = "Loose Coupling: Services operate independently, reducing interdependencies."
p.level = 2
p = content.add_paragraph()
p.text = "Independent Deployability: Each service can be updated without affecting others."
p.level = 2
p = content.add_paragraph()
p.text = "Technology Diversity: Teams can choose optimal tools for specific services."
p.level = 2
p = content.add_paragraph()
p.text = "Benefits:"
p.level = 1
p = content.add_paragraph()
p.text = "Faster development and deployment cycles."
p.level = 2
p = content.add_paragraph()
p.text = "Improved scalability for individual components."
p.level = 2
p = content.add_paragraph()
p.text = "Enhanced fault isolation, ensuring system resilience."
p.level = 2
# Add an icon (simulated as a shape)
icon = slide.shapes.add_shape(MSO_SHAPE.GEAR_6_POINTS, Inches(8), Inches(0.5), Inches(1), Inches(1))
icon.fill.solid()
icon.fill.fore_color.rgb = RGBColor(0, 128, 0)
icon.line.color.rgb = RGBColor(0, 0, 0)

# Slide 4: Proposed Microservices Architecture (with Diagram)
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Proposed Microservices Architecture"
content = slide.placeholders[1].text_frame
content.text = "Structure"
p = content.add_paragraph()
p.text = "Main Services:"
p.level = 1
p = content.add_paragraph()
p.text = "Booking Service: Manages customer bookings and availability."
p.level = 2
p = content.add_paragraph()
p.text = "Scheduling Service: Assigns cleaners and optimises schedules."
p.level = 2
p = content.add_paragraph()
p.text = "Invoicing Service: Handles billing and payment processing."
p.level = 2
p = content.add_paragraph()
p.text = "Customer Management Service: Stores and manages customer data."
p.level = 2
p = content.add_paragraph()
p.text = "Responsibilities:"
p.level = 1
p = content.add_paragraph()
p.text = "Each service handles a specific business function, aligned with domain-driven design."
p.level = 2
p = content.add_paragraph()
p.text = "Communication Patterns:"
p.level = 1
p = content.add_paragraph()
p.text = "API Gateway: Central entry point for client requests, routing to appropriate services."
p.level = 2
p = content.add_paragraph()
p.text = "Event-Driven Communication: Asynchronous messaging for updates."
p.level = 2
p = content.add_paragraph()
p.text = "Database per Service: Each service has its own database to ensure decoupling."
p.level = 2
# Add a diagram
# API Gateway
api_gateway = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(1.5), Inches(2), Inches(0.5))
api_gateway.fill.solid()
api_gateway.fill.fore_color.rgb = RGBColor(0, 105, 148)
api_gateway_text = api_gateway.text_frame
api_gateway_text.text = "API Gateway"
# Services
booking_service = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2.5), Inches(1.5), Inches(0.5))
booking_service.fill.solid()
booking_service.fill.fore_color.rgb = RGBColor(46, 204, 113)
booking_service_text = booking_service.text_frame
booking_service_text.text = "Booking"
scheduling_service = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(2.5), Inches(1.5), Inches(0.5))
scheduling_service.fill.solid()
scheduling_service.fill.fore_color.rgb = RGBColor(46, 204, 113)
scheduling_service_text = scheduling_service.text_frame
scheduling_service_text.text = "Scheduling"
invoicing_service = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(2.5), Inches(1.5), Inches(0.5))
invoicing_service.fill.solid()
invoicing_service.fill.fore_color.rgb = RGBColor(46, 204, 113)
invoicing_service_text = invoicing_service.text_frame
invoicing_service_text.text = "Invoicing"
# Arrows to show communication
arrow1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.75), Inches(2), Inches(0.5), Inches(0.5))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = RGBColor(0, 0, 0)
arrow2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(3.25), Inches(2), Inches(0.5), Inches(0.5))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = RGBColor(0, 0, 0)
arrow3 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.25), Inches(2), Inches(0.5), Inches(0.5))
arrow3.fill.solid()
arrow3.fill.fore_color.rgb = RGBColor(0, 0, 0)

# Slide 5: Comparison - Monolithic vs. Microservices
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Comparison - Monolithic vs. Microservices"
table = slide.shapes.add_table(6, 3, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table
table.columns[0].width = Inches(2)
table.columns[1].width = Inches(3.5)
table.columns[2].width = Inches(3.5)
table.cell(0, 0).text = "Aspect"
table.cell(0, 1).text = "Monolithic Architecture"
table.cell(0, 2).text = "Microservices Architecture"
table.cell(1, 0).text = "Development Agility"
table.cell(1, 1).text = "Slow due to complex, unified codebase"
table.cell(1, 2).text = "Faster with independent, smaller codebases"
table.cell(2, 0).text = "Scalability"
table.cell(2, 1).text = "Scales as a single unit, resource-heavy"
table.cell(2, 2).text = "Scales specific services as needed"
table.cell(3, 0).text = "Fault Isolation"
table.cell(3, 1).text = "Single failure impacts entire system"
table.cell(3, 2).text = "Failures confined to individual services"
table.cell(4, 0).text = "Deployment"
table.cell(4, 1).text = "Full system redeployment, frequent downtime"
table.cell(4, 2).text = "Independent deployments, minimal downtime"
table.cell(5, 0).text = "Business Alignment"
table.cell(5, 1).text = "Rigid, slows response to market changes"
table.cell(5, 2).text = "Flexible, supports rapid feature delivery"
for row in table.rows:
    for cell in row.cells:
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        if row == table.rows[0]:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 105, 148)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
# Add an icon (simulated as a shape)
icon = slide.shapes.add_shape(MSO_SHAPE.PIE, Inches(8), Inches(0.5), Inches(1), Inches(1))
icon.fill.solid()
icon.fill.fore_color.rgb = RGBColor(0, 105, 148)
icon.line.color.rgb = RGBColor(0, 0, 0)

# Slide 6: Recommendations and Business Benefits
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Recommendations and Business Benefits"
content = slide.placeholders[1].text_frame
content.text = "Action Plan"
p = content.add_paragraph()
p.text = "Adopt Microservices Architecture:"
p.level = 1
p = content.add_paragraph()
p.text = "Decompose monolithic system into proposed services."
p.level = 2
p = content.add_paragraph()
p.text = "Implement API gateway and event-driven communication."
p.level = 2
p = content.add_paragraph()
p.text = "Business Benefits:"
p.level = 1
p = content.add_paragraph()
p.text = "Adaptability: Quickly roll out new features (e.g., subscription-based cleaning plans)."
p.level = 2
p = content.add_paragraph()
p.text = "Scalability: Efficiently handle peak demand (e.g., holiday season bookings)."
p.level = 2
p = content.add_paragraph()
p.text = "Resilience: Minimise downtime, ensuring consistent customer experience."
p.level = 2
p = content.add_paragraph()
p.text = "Alignment with Objectives: Supports growth by enabling faster market response."
p.level = 2
p = content.add_paragraph()
p.text = "Next Steps:"
p.level = 1
p = content.add_paragraph()
p.text = "Pilot a single service (e.g., Booking Service) to validate approach."
p.level = 2
p = content.add_paragraph()
p.text = "Train teams on microservices best practices and DevOps tools."
p.level = 2
# Add an icon (simulated as a shape)
icon = slide.shapes.add_shape(MSO_SHAPE.LIGHTNING_BOLT, Inches(8), Inches(0.5), Inches(1), Inches(1))
icon.fill.solid()
icon.fill.fore_color.rgb = RGBColor(255, 215, 0)
icon.line.color.rgb = RGBColor(0, 0, 0)

# Slide 7: Conclusion
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Conclusion"
content = slide.placeholders[1].text_frame
content.text = "Summary"
p = content.add_paragraph()
p.text = "Transitioning to a microservices architecture addresses the limitations of the monolithic system."
p.level = 1
p = content.add_paragraph()
p.text = "Enables better adaptability, scalability, and resilience."
p.level = 1
p = content.add_paragraph()
p.text = "Aligns with business goals of improving customer satisfaction and operational efficiency."
p.level = 1
p = content.add_paragraph()
p.text = "Call to Action: Approve pilot project to begin transformation."
p.level = 1
# Add an icon (simulated as a shape)
icon = slide.shapes.add_shape(MSO_SHAPE.CHECKMARK, Inches(8), Inches(0.5), Inches(1), Inches(1))
icon.fill.solid()
icon.fill.fore_color.rgb = RGBColor(0, 128, 0)
icon.line.color.rgb = RGBColor(0, 0, 0)

# Save the presentation
prs.save("Microservices_Architecture_Presentation_With_Visuals.pptx")