from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Helper function to add a slide with title and bullet points
def add_slide(prs, title, bullet_points):
    slide_layout = prs.slide_layouts[1]  # Bullet slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
    
    # Add bullet points
    content = slide.placeholders[1]
    text_frame = content.text_frame
    for point in bullet_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.LEFT

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Revised Data Ingestion Architecture for ABM Facility"
title.text_frame.paragraphs[0].font.size = Pt(40)
subtitle = slide.placeholders[1]
subtitle.text = "Enhancing Operational Efficiency and Scalability\nPresented by: Your Name\nDate: July 21, 2025\nABM Facility: Transforming Cleaning Operations"
subtitle.text_frame.paragraphs[0].font.size = Pt(20)

# Slide 2: Executive Summary
add_slide(prs, "Executive Summary", [
    "Revised pipeline improves efficiency, scalability, and security.",
    "Integrates new data sources for real-time analytics.",
    "Enhances job scheduling, inventory tracking, and quality control.",
    "Positions ABM Facility for operational excellence and proactive decision-making."
])

# Slide 3: Challenges and Objectives
add_slide(prs, "Challenges and Objectives", [
    "Challenges:",
    "- Manual scheduling caused duplicate tasks and stock issues.",
    "- Lack of real-time data from mobile forms and IoT sensors.",
    "- Inconsistent quality feedback and equipment downtime.",
    "Objectives:",
    "- Build a scalable data ingestion pipeline.",
    "- Enable real-time analytics and proactive decisions.",
    "- Improve consistency in service delivery."
])

# Slide 4: Revised Architecture Overview
add_slide(prs, "Revised Architecture Overview", [
    "Components:",
    "- Data Sources: Microsoft Forms (Excel), SQL Server, IoT Sensors, REST APIs.",
    "- Landing Layer: Azure Blob Storage, Event Grid, Data Gateway.",
    "- Ingestion Layer: Azure Functions for validation, Azure Data Factory for ETL workflows.",
    "Benefits:",
    "- Enhanced reliability, scalability, and security.",
    "- Real-time data processing for better decision-making."
])

# Slide 5: Key Features
add_slide(prs, "Key Features", [
    "Data Validation: Azure Functions ensure schema consistency.",
    "Duplicate Detection: SQL queries eliminate redundant job entries.",
    "Security: Azure Virtual Networks, private endpoints, and PII masking.",
    "Monitoring: Azure Security Centre and Sentinel for real-time threat detection.",
    "Logging: Azure Monitor retains logs for 12 months for audits."
])

# Slide 6: Results and Impact
add_slide(prs, "Results and Impact", [
    "Achievements:",
    "- 40% reduction in duplicate/missed tasks.",
    "- 30% less equipment downtime via IoT predictive maintenance.",
    "- 15% reduction in unnecessary inventory orders.",
    "Case Study:",
    "- Predictive maintenance alerts improved client satisfaction and cut repair costs.",
    "- Power BI dashboards optimized procurement processes."
])

# Slide 7: Recommendations and Future Steps
add_slide(prs, "Recommendations and Future Steps", [
    "Recommendations:",
    "- Integrate automated anomaly detection in Azure Functions.",
    "- Expand IoT sensor deployment for predictive maintenance.",
    "- Adopt Apache Kafka or Azure Event Hubs for high-volume data streams.",
    "- Enhance GDPR compliance with real-time data integration.",
    "Future Vision:",
    "- Leverage Kubernetes for scalability and serverless computing for efficiency.",
    "- Ensure adaptability through continuous monitoring and feedback."
])

# Save the presentation
prs.save("ABMFacilityDataIngestion.pptx")
print("Presentation saved as ABMFacilityDataIngestion.pptx")